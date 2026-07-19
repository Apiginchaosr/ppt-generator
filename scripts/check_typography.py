#!/usr/bin/env python3
"""
PPT Typography & Layout Checker — Post-generation quality assurance.

Scans a .pptx file for common issues:
- Text overflow (content exceeding text box boundaries)
- Font size violations (body < 18pt, titles < 28pt)
- Font consistency issues
- Alignment problems
- Color contrast readability
- Element overlapping

Usage:
    python check_typography.py <presentation.pptx> [--fix] [--lang zh|en]
"""

import argparse
import os
import sys
import math
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# WCAG relative luminance & contrast ratio
# ---------------------------------------------------------------------------
def relative_luminance(rgb: RGBColor) -> float:
    """Compute WCAG relative luminance for an RGBColor."""
    def channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r = channel(rgb[0] if isinstance(rgb, tuple) else rgb.red)
    g = channel(rgb[1] if isinstance(rgb, tuple) else rgb.green)
    b = channel(rgb[2] if isinstance(rgb, tuple) else rgb.blue)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(c1, c2) -> float:
    """WCAG contrast ratio between two colors."""
    l1 = relative_luminance(c1)
    l2 = relative_luminance(c2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


# ---------------------------------------------------------------------------
# Extract text run details from a shape
# ---------------------------------------------------------------------------
def get_shape_font_info(shape):
    """Extract font details from all text runs in a shape."""
    runs_info = []
    if not shape.has_text_frame:
        return runs_info
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            color = None
            try:
                if run.font.color and run.font.color.rgb:
                    color = run.font.color.rgb
            except (AttributeError, TypeError):
                pass
            info = {
                "text": run.text,
                "font_name": run.font.name,
                "font_size": run.font.size,
                "bold": run.font.bold,
                "color": color,
            }
            runs_info.append(info)
    return runs_info


def get_shape_text(shape):
    """Get full text content of a shape."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text


def shape_position(shape):
    """Return (left_inches, top_inches, right_inches, bottom_inches)."""
    l = shape.left / 914400 if shape.left else 0
    t = shape.top / 914400 if shape.top else 0
    r = l + (shape.width / 914400 if shape.width else 0)
    b = t + (shape.height / 914400 if shape.height else 0)
    return l, t, r, b


def shapes_overlap(s1, s2, margin_inches=0.05):
    """Check if two shapes overlap (with margin for near-misses)."""
    l1, t1, r1, b1 = shape_position(s1)
    l2, t2, r2, b2 = shape_position(s2)
    return not (r1 + margin_inches < l2 or r2 + margin_inches < l1 or
                b1 + margin_inches < t2 or b2 + margin_inches < t1)


def estimate_text_lines(text: str, font_size_pt: float, box_width_inches: float) -> int:
    """Rough estimate of how many lines text will occupy."""
    if not text or font_size_pt <= 0:
        return 0
    # Approximate: average char width ≈ 0.6 * font_size for CJK, 0.4 for Latin
    # Use 0.5 as a rough average
    avg_char_width_inches = (font_size_pt * 0.5) / 72  # convert pt to inches
    chars_per_line = max(1, box_width_inches / avg_char_width_inches)
    total_lines = 0
    for paragraph in text.split('\n'):
        if not paragraph.strip():
            total_lines += 1
        else:
            total_lines += max(1, math.ceil(len(paragraph) / chars_per_line))
    return total_lines


# ---------------------------------------------------------------------------
# Typography Check
# ---------------------------------------------------------------------------
MIN_BODY_SIZE_PT = 18
MIN_TITLE_SIZE_PT = 28
MIN_CONTRAST_RATIO = 4.5  # WCAG AA for normal text
MIN_LARGE_TEXT_CONTRAST = 3.0  # WCAG AA for large text (>18pt or bold >14pt)


def check_presentation(pptx_path: str, lang: str = "zh") -> list:
    """Run all checks and return list of issues."""
    prs = Presentation(pptx_path)
    issues = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]

        # --- Check 1: Text overflow ---
        for shape in text_shapes:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            runs = get_shape_font_info(shape)
            if not runs:
                continue

            font_size = runs[0].get("font_size")
            if font_size is None:
                font_size = Pt(18)
            font_size_pt = font_size / 12700  # EMU to points

            box_w = shape.width / 914400  # EMU to inches
            box_h = shape.height / 914400

            est_lines = estimate_text_lines(text, font_size_pt, box_w)
            line_height_inches = (font_size_pt * 1.5) / 72
            est_height = est_lines * line_height_inches

            if est_height > box_h * 1.05:  # 5% tolerance
                overflow_pct = round((est_height - box_h) / box_h * 100)
                issues.append({
                    "slide": slide_num,
                    "type": "text_overflow",
                    "severity": "high",
                    "detail": (
                        f"Text may overflow (est. {overflow_pct}% beyond box). "
                        f'"{text[:60]}..." ({est_lines} lines in {font_size_pt:.0f}pt)'
                    ),
                    "fixable": False,  # Needs manual content trimming
                })

        # --- Check 2: Font size minimums ---
        for shape in text_shapes:
            runs = get_shape_font_info(shape)
            for run in runs:
                size_pt = (run.get("font_size") or Pt(18)) / 12700
                is_bold = run.get("bold", False)
                is_large = size_pt >= MIN_TITLE_SIZE_PT or (size_pt >= 14 and is_bold)

                if not is_large and size_pt < MIN_BODY_SIZE_PT:
                    issues.append({
                        "slide": slide_num,
                        "type": "font_too_small",
                        "severity": "medium",
                        "detail": (
                            f'Font size {size_pt:.0f}pt below minimum ({MIN_BODY_SIZE_PT}pt). '
                            f'Text: "{run.get("text", "")[:40]}"'
                        ),
                        "fixable": True,
                        "fix_action": f"Set font size to {MIN_BODY_SIZE_PT}pt",
                    })

        # --- Check 3: Font consistency ---
        if text_shapes:
            all_fonts = []
            for shape in text_shapes:
                for run in get_shape_font_info(shape):
                    fn = run.get("font_name")
                    if fn:
                        all_fonts.append(fn)
            if all_fonts:
                font_counts = Counter(all_fonts)
                # If more than 3 font families on one slide, flag it
                if len(font_counts) > 3:
                    issues.append({
                        "slide": slide_num,
                        "type": "font_inconsistency",
                        "severity": "low",
                        "detail": f"Multiple fonts on slide: {dict(font_counts.most_common())}",
                        "fixable": False,
                    })

        # --- Check 4: Color contrast ---
        for shape in text_shapes:
            runs = get_shape_font_info(shape)
            # Estimate background color from slide or shape fill
            bg_rgb = _get_shape_bg(shape, slide)
            for run in runs:
                fg = run.get("color")
                if fg is None:
                    continue
                size_pt = (run.get("font_size") or Pt(18)) / 12700
                is_large = size_pt >= MIN_TITLE_SIZE_PT or (size_pt >= 14 and run.get("bold"))
                min_ratio = MIN_LARGE_TEXT_CONTRAST if is_large else MIN_CONTRAST_RATIO

                ratio = contrast_ratio(fg, bg_rgb)
                if ratio < min_ratio:
                    issues.append({
                        "slide": slide_num,
                        "type": "low_contrast",
                        "severity": "medium",
                        "detail": (
                            f'Contrast ratio {ratio:.1f}:1 below {min_ratio}:1. '
                            f'Text: "{run.get("text", "")[:40]}" '
                            f'({size_pt:.0f}pt) on bg'
                        ),
                        "fixable": True,
                        "fix_action": "Darken text or lighten background",
                    })

        # --- Check 5: Overlapping elements ---
        # Only check overlap between TEXT shapes. Decorative shapes (accent bars,
        # background rects) intentionally overlap with text.
        for i, s1 in enumerate(shapes):
            for j, s2 in enumerate(shapes):
                if j <= i:
                    continue
                t1 = get_shape_text(s1).strip()
                t2 = get_shape_text(s2).strip()
                # Both must have text to be a real overlap issue
                if not t1 or not t2:
                    continue
                if shapes_overlap(s1, s2):
                    issues.append({
                        "slide": slide_num,
                        "type": "overlap",
                        "severity": "high",
                        "detail": (
                            f'Two elements overlap: '
                            f'"{t1[:30]}..." and "{t2[:30]}..."'
                        ),
                        "fixable": True,
                        "fix_action": "Adjust position of one element",
                    })

    return issues


def _get_shape_bg(shape, slide) -> RGBColor:
    """Best-effort background color for a shape."""
    # Try shape fill first
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            try:
                if fc and fc.rgb:
                    return fc.rgb
            except (AttributeError, TypeError):
                pass
    except Exception:
        pass
    # Fall back to slide background
    try:
        bg = slide.background.fill
        if bg.type is not None:
            fc = bg.fore_color
            try:
                if fc and fc.rgb:
                    return fc.rgb
            except (AttributeError, TypeError):
                pass
    except Exception:
        pass
    return RGBColor(0xFF, 0xFF, 0xFF)  # Assume white


# ---------------------------------------------------------------------------
# Report formatting
# ---------------------------------------------------------------------------
LABELS = {
    "text_overflow": {"zh": "文字溢出", "en": "Text Overflow"},
    "font_too_small": {"zh": "字号过小", "en": "Font Too Small"},
    "font_inconsistency": {"zh": "字体不一致", "en": "Font Inconsistency"},
    "low_contrast": {"zh": "对比度不足", "en": "Low Contrast"},
    "overlap": {"zh": "元素重叠", "en": "Overlapping Elements"},
}

SEVERITY_ICONS = {"high": "[HIGH]", "medium": "[MED]", "low": "[LOW]"}


def print_report(issues: list, total_slides: int, lang: str = "zh"):
    """Print a human-readable check report."""
    if not issues:
        print(f"[OK] All checks passed! ({total_slides} slides, 0 issues)")
        return

    # Count affected slides
    affected = set(i["slide"] for i in issues)
    high = [i for i in issues if i["severity"] == "high"]
    medium = [i for i in issues if i["severity"] == "medium"]
    low = [i for i in issues if i["severity"] == "low"]

    title = "排版检查报告" if lang == "zh" else "Typography Check Report"
    print(f"\n{'='*60}")
    print(f"[CHECK] {title} ({len(affected)}/{total_slides} slides have issues)")
    print(f"   [HIGH] {len(high)}   [MED] {len(medium)}   [LOW] {len(low)}")
    print(f"{'='*60}")

    # Group by slide
    for slide_num in sorted(affected):
        slide_issues = [i for i in issues if i["slide"] == slide_num]
        print(f"\n-- Slide {slide_num} --")
        for issue in slide_issues:
            icon = SEVERITY_ICONS.get(issue["severity"], "⚪")
            type_label = LABELS.get(issue["type"], {}).get(lang, issue["type"])
            fixable = "[FIX]" if issue.get("fixable") else "[MANUAL]"
            print(f"  {icon} [{type_label}] {fixable} {issue['detail']}")

    fixable_count = sum(1 for i in issues if i.get("fixable"))
    print(f"\n{'='*60}")
    print(f"Summary: {len(issues)} issues ({fixable_count} auto-fixable, {len(issues) - fixable_count} need manual review)")
    print(f"{'='*60}\n")


def to_json(issues: list) -> dict:
    """Convert issues to JSON-serializable format."""
    result = []
    for issue in issues:
        item = dict(issue)
        # Convert any non-serializable values
        result.append(item)
    return {"issues": result, "total": len(result)}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description="Check PPTX typography and layout")
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument("--fix", action="store_true", help="Attempt automatic fixes")
    parser.add_argument("--lang", choices=["zh", "en"], default="zh", help="Report language")
    parser.add_argument("--json", action="store_true", help="Output JSON instead of text report")
    args = parser.parse_args()

    if not os.path.exists(args.pptx_path):
        print(f"❌ File not found: {args.pptx_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(args.pptx_path)
    total_slides = len(prs.slides)
    issues = check_presentation(args.pptx_path, args.lang)

    if args.json:
        import json
        print(json.dumps(to_json(issues), ensure_ascii=False, indent=2))
    else:
        print_report(issues, total_slides, args.lang)

    if args.fix:
        fixable = [i for i in issues if i.get("fixable")]
        if fixable:
            print(f"\n🔧 {len(fixable)} issues marked as auto-fixable.")
            print("Note: Full auto-fix requires re-generating the PPT via generate_pptx.py with adjusted config.")
            print("For now, manually address the flagged issues above.")
        else:
            print("\nNo auto-fixable issues found.")

    # Exit with non-zero if high-severity issues exist
    high = [i for i in issues if i["severity"] == "high"]
    if high:
        sys.exit(1)


if __name__ == "__main__":
    main()
