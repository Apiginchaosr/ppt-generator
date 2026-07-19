#!/usr/bin/env python3
"""
PPT Generator — Core PPTX generation engine for the ppt-generator skill.

Takes a JSON config file describing slides, design, and content,
produces a .pptx file ready for use.

Usage:
    python generate_pptx.py --config config.json --output presentation.pptx
"""

import argparse
import json
import os
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# Default design presets by scenario
# ---------------------------------------------------------------------------
SCENARIO_PRESETS = {
    "business": {
        "palette": {
            "primary": "#1a3a5c",
            "secondary": "#f0f4f8",
            "accent": "#e8a817",
            "text": "#333333",
            "background": "#ffffff",
        },
        "fonts": {
            "title": "Microsoft YaHei",
            "body": "Microsoft YaHei",
            "title_size_pt": 32,
            "body_size_pt": 18,
            "subtitle_size_pt": 20,
        },
    },
    "academic": {
        "palette": {
            "primary": "#2c3e50",
            "secondary": "#ecf0f1",
            "accent": "#8e44ad",
            "text": "#2c3e50",
            "background": "#ffffff",
        },
        "fonts": {
            "title": "SimHei",
            "body": "Microsoft YaHei",
            "title_size_pt": 30,
            "body_size_pt": 18,
            "subtitle_size_pt": 20,
        },
    },
    "tech": {
        "palette": {
            "primary": "#00d4aa",
            "secondary": "#1e1e2e",
            "accent": "#f38ba8",
            "text": "#cdd6f4",
            "background": "#1e1e2e",
        },
        "fonts": {
            "title": "Consolas",
            "body": "Microsoft YaHei",
            "title_size_pt": 32,
            "body_size_pt": 18,
            "subtitle_size_pt": 20,
        },
    },
    "creative": {
        "palette": {
            "primary": "#ff6b6b",
            "secondary": "#feca57",
            "accent": "#48dbfb",
            "text": "#2d3436",
            "background": "#ffffff",
        },
        "fonts": {
            "title": "Microsoft YaHei",
            "body": "Microsoft YaHei",
            "title_size_pt": 36,
            "body_size_pt": 20,
            "subtitle_size_pt": 24,
        },
    },
    "general": {
        "palette": {
            "primary": "#4a4a4a",
            "secondary": "#f5f5f5",
            "accent": "#0078d4",
            "text": "#333333",
            "background": "#ffffff",
        },
        "fonts": {
            "title": "Microsoft YaHei",
            "body": "Microsoft YaHei",
            "title_size_pt": 32,
            "body_size_pt": 18,
            "subtitle_size_pt": 20,
        },
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#1a3a5c' to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def emu(inches: float) -> Emu:
    return Inches(inches)


def set_slide_bg(slide, color_hex: str):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_name: str = "Microsoft YaHei",
    font_size_pt: float = 18,
    bold: bool = False,
    color_hex: str = "#333333",
    alignment=PP_ALIGN.LEFT,
    line_spacing: float = 1.2,
):
    """Add a text box to a slide and return (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size_pt)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.alignment = alignment
    p.line_spacing = Pt(font_size_pt * line_spacing)
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn("a:altLang"), "zh-CN")
    return txBox, tf


def add_multiline_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list,
    font_name: str = "Microsoft YaHei",
    font_size_pt: float = 18,
    color_hex: str = "#333333",
    bullet: bool = False,
    line_spacing: float = 1.5,
):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = hex_to_rgb(color_hex)
        p.line_spacing = Pt(font_size_pt * line_spacing)
        if bullet:
            p.level = 0
    return txBox, tf


def add_image_safe(slide, image_path: str, left, top, width, height):
    """Add image if file exists, otherwise add a placeholder rectangle."""
    if image_path and os.path.exists(image_path):
        return slide.shapes.add_picture(image_path, emu(left), emu(top), emu(width), emu(height))
    else:
        # Placeholder rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(width), emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        # Add placeholder text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[图片]"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER
        return shape


def add_shape_bg(slide, left, top, width, height, color_hex: str):
    """Add a filled rectangle (useful for colored bands, accent bars)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(width), emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.background()  # No border
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_cover(slide, slide_data: dict, design: dict):
    """Build a cover/title slide."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["background"])

    # Accent bar at top
    add_shape_bg(slide, 0, 0, 13.333, 0.15, palette["primary"])

    # Title
    title = slide_data.get("title", "Title")
    add_textbox(
        slide,
        left=1.5,
        top=2.2,
        width=10.3,
        height=1.5,
        text=title,
        font_name=fonts["title"],
        font_size_pt=fonts.get("title_size_pt", 32),
        bold=True,
        color_hex=palette["primary"],
        alignment=PP_ALIGN.LEFT,
    )

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            left=1.5,
            top=3.8,
            width=10.3,
            height=0.8,
            text=subtitle,
            font_name=fonts["body"],
            font_size_pt=fonts.get("subtitle_size_pt", 20),
            bold=False,
            color_hex=palette.get("text", "#666666"),
            alignment=PP_ALIGN.LEFT,
        )

    # Accent bar at bottom
    add_shape_bg(slide, 0, 7.35, 13.333, 0.15, palette["accent"])


def build_toc(slide, slide_data: dict, design: dict):
    """Build a table of contents slide."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["background"])

    title = slide_data.get("title", "目录")
    add_textbox(
        slide,
        left=1.0,
        top=0.6,
        width=5.0,
        height=0.8,
        text=title,
        font_name=fonts["title"],
        font_size_pt=fonts.get("title_size_pt", 32),
        bold=True,
        color_hex=palette["primary"],
    )

    # Accent line under title
    add_shape_bg(slide, 1.0, 1.45, 2.0, 0.05, palette["accent"])

    items = slide_data.get("items", [])
    y_start = 2.0
    for i, item in enumerate(items):
        num = f"0{i+1}" if i + 1 < 10 else str(i + 1)
        # Number
        add_textbox(
            slide,
            left=1.5,
            top=y_start + i * 0.7,
            width=0.6,
            height=0.5,
            text=num,
            font_name=fonts["title"],
            font_size_pt=24,
            bold=True,
            color_hex=palette["accent"],
        )
        # Item text
        add_textbox(
            slide,
            left=2.2,
            top=y_start + i * 0.7,
            width=8.0,
            height=0.5,
            text=item,
            font_name=fonts["body"],
            font_size_pt=20,
            bold=False,
            color_hex=palette["text"],
        )


def build_content(slide, slide_data: dict, design: dict):
    """Build a general content slide with configurable layout."""
    palette = design["palette"]
    fonts = design["fonts"]
    layout = slide_data.get("layout", "text_only")
    set_slide_bg(slide, palette["background"])

    title = slide_data.get("title", "")
    body = slide_data.get("body", [])
    image_data = slide_data.get("image")

    # Title bar
    add_shape_bg(slide, 0, 0, 13.333, 0.9, palette["primary"])
    add_textbox(
        slide,
        left=0.8,
        top=0.1,
        width=11.7,
        height=0.7,
        text=title,
        font_name=fonts["title"],
        font_size_pt=28,
        bold=True,
        color_hex="#ffffff",
    )

    # Body area varies by layout
    layouts = {
        "text_only": _layout_text_only,
        "text_left_image_right": _layout_text_left_image_right,
        "image_left_text_right": _layout_image_left_text_right,
        "text_top_image_bottom": _layout_text_top_image_bottom,
        "two_column": _layout_two_column,
        "three_column": _layout_three_column,
    }
    builder = layouts.get(layout, _layout_text_only)
    builder(slide, body, image_data, design, palette, fonts)


def _layout_text_only(slide, body, image_data, design, palette, fonts):
    add_multiline_textbox(
        slide,
        left=1.0,
        top=1.3,
        width=11.3,
        height=5.8,
        lines=body if isinstance(body, list) else [body],
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )


def _layout_text_left_image_right(slide, body, image_data, design, palette, fonts):
    add_multiline_textbox(
        slide,
        left=0.8,
        top=1.3,
        width=6.5,
        height=5.5,
        lines=body if isinstance(body, list) else [body],
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )
    if image_data:
        path = image_data.get("path", "")
        add_image_safe(slide, path, left=7.8, top=1.5, width=4.8, height=4.5)


def _layout_image_left_text_right(slide, body, image_data, design, palette, fonts):
    if image_data:
        path = image_data.get("path", "")
        add_image_safe(slide, path, left=0.5, top=1.5, width=5.0, height=4.5)
    add_multiline_textbox(
        slide,
        left=6.0,
        top=1.3,
        width=6.8,
        height=5.5,
        lines=body if isinstance(body, list) else [body],
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )


def _layout_text_top_image_bottom(slide, body, image_data, design, palette, fonts):
    add_multiline_textbox(
        slide,
        left=0.8,
        top=1.3,
        width=11.5,
        height=2.5,
        lines=body if isinstance(body, list) else [body],
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )
    if image_data:
        path = image_data.get("path", "")
        add_image_safe(slide, path, left=1.5, top=4.0, width=10.3, height=3.0)


def _layout_two_column(slide, body, image_data, design, palette, fonts):
    """body should be {'left': [...], 'right': [...]} or a flat list split in half."""
    if isinstance(body, dict):
        left_items = body.get("left", [])
        right_items = body.get("right", [])
    else:
        mid = (len(body) + 1) // 2
        left_items = body[:mid]
        right_items = body[mid:]

    add_multiline_textbox(
        slide,
        left=0.5,
        top=1.3,
        width=5.8,
        height=5.5,
        lines=left_items,
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )
    # Vertical divider
    add_shape_bg(slide, 6.6, 1.3, 0.02, 5.5, palette["accent"])
    add_multiline_textbox(
        slide,
        left=7.0,
        top=1.3,
        width=5.8,
        height=5.5,
        lines=right_items,
        font_name=fonts["body"],
        font_size_pt=fonts.get("body_size_pt", 18),
        color_hex=palette["text"],
        bullet=True,
    )


def _layout_three_column(slide, body, image_data, design, palette, fonts):
    """body should be {'col1': [...], 'col2': [...], 'col3': [...]}."""
    cols = [body.get(f"col{i}", []) if isinstance(body, dict) else [] for i in range(1, 4)]
    if not any(cols) and isinstance(body, list):
        # Split evenly
        chunk = (len(body) + 2) // 3
        cols = [body[i:i+chunk] for i in range(0, len(body), chunk)]

    col_width = 4.0
    for idx, col_items in enumerate(cols[:3]):
        x = 0.5 + idx * (col_width + 0.3)
        add_multiline_textbox(
            slide,
            left=x,
            top=1.3,
            width=col_width,
            height=5.5,
            lines=col_items,
            font_name=fonts["body"],
            font_size_pt=fonts.get("body_size_pt", 16),
            color_hex=palette["text"],
            bullet=True,
        )
        if idx < 2:
            add_shape_bg(slide, x + col_width + 0.1, 1.3, 0.02, 5.5, palette["secondary"])


def build_chart(slide, slide_data: dict, design: dict):
    """Build a slide with an embedded chart."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["background"])

    title = slide_data.get("title", "Chart")
    chart_type_str = slide_data.get("chart_type", "bar")
    data = slide_data.get("data", {})

    # Title bar
    add_shape_bg(slide, 0, 0, 13.333, 0.9, palette["primary"])
    add_textbox(
        slide, left=0.8, top=0.1, width=11.7, height=0.7,
        text=title, font_name=fonts["title"], font_size_pt=28,
        bold=True, color_hex="#ffffff",
    )

    # Map chart type
    chart_types = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart_type = chart_types.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    chart_data = CategoryChartData()
    categories = data.get("categories", [])
    chart_data.categories = categories

    series_dict = data.get("series", {})
    for series_name, values in series_dict.items():
        # Ensure values length matches categories
        padded = values + [0] * (len(categories) - len(values))
        chart_data.add_series(series_name, padded[:len(categories)])

    # Add chart shape
    chart_frame = slide.shapes.add_chart(
        chart_type, emu(1.0), emu(1.3), emu(11.3), emu(5.5), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = len(series_dict) > 1

    # Style the chart
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)

    # Color the series
    series_colors = [palette["primary"], palette["accent"], palette["secondary"], "#888888"]
    for i, series in enumerate(chart.series):
        if i < len(series_colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = hex_to_rgb(series_colors[i])


def build_table(slide, slide_data: dict, design: dict):
    """Build a slide with a formatted table."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["background"])

    title = slide_data.get("title", "Table")
    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])

    # Title bar
    add_shape_bg(slide, 0, 0, 13.333, 0.9, palette["primary"])
    add_textbox(
        slide, left=0.8, top=0.1, width=11.7, height=0.7,
        text=title, font_name=fonts["title"], font_size_pt=28,
        bold=True, color_hex="#ffffff",
    )

    if not headers or not rows:
        add_textbox(slide, 1.0, 1.5, 11.3, 1.0, "[No table data provided]",
                    color_hex=palette["text"])
        return

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_width = 11.3
    table_height = min(5.5, 0.5 * n_rows)
    col_width = table_width / n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, emu(1.0), emu(1.3), emu(table_width), emu(table_height)
    )
    table = table_shape.table

    # Set column widths
    for ci in range(n_cols):
        table.columns[ci].width = emu(col_width)

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(palette["primary"])
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = fonts["body"]
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(palette["secondary"])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = hex_to_rgb(palette["text"])
                p.font.name = fonts["body"]
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT


def build_image_slide(slide, slide_data: dict, design: dict):
    """Full-slide image with optional overlay text."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, "#000000")

    image_data = slide_data.get("image", {})
    path = image_data.get("path", "")
    if path and os.path.exists(path):
        # Full slide image
        slide.shapes.add_picture(path, emu(0), emu(0), emu(13.333), emu(7.5))

    # Overlay text if provided
    overlay_title = slide_data.get("title", "")
    overlay_sub = slide_data.get("subtitle", "")
    if overlay_title:
        # Semi-transparent overlay bar at bottom
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, emu(0), emu(5.5), emu(13.333), emu(2.0)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0, 0, 0)
        bar.line.fill.background()
        # Set transparency via XML
        solidFill = bar.fill._fill
        srgb = solidFill.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = srgb.makeelement(qn("a:alpha"), {"val": "50000"})  # 50% opacity
            srgb.append(alpha)

        add_textbox(slide, 1.0, 5.7, 11.3, 0.7, overlay_title,
                    font_name=fonts["title"], font_size_pt=30,
                    bold=True, color_hex="#ffffff")
        if overlay_sub:
            add_textbox(slide, 1.0, 6.4, 11.3, 0.5, overlay_sub,
                        font_name=fonts["body"], font_size_pt=18,
                        color_hex="#cccccc")


def build_diagram(slide, slide_data: dict, design: dict):
    """Build a diagram slide (pre-rendered Mermaid or other diagram as image)."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["background"])

    title = slide_data.get("title", "Diagram")
    image_data = slide_data.get("image", {})

    add_shape_bg(slide, 0, 0, 13.333, 0.9, palette["primary"])
    add_textbox(
        slide, left=0.8, top=0.1, width=11.7, height=0.7,
        text=title, font_name=fonts["title"], font_size_pt=28,
        bold=True, color_hex="#ffffff",
    )

    path = image_data.get("path", "")
    if path and os.path.exists(path):
        add_image_safe(slide, path, left=1.5, top=1.2, width=10.3, height=5.5)
    else:
        add_multiline_textbox(
            slide, left=1.5, top=1.5, width=10.3, height=4.0,
            lines=["[Diagram placeholder]", f"Mermaid: {image_data.get('mermaid', 'N/A')}"],
            font_name=fonts["body"], font_size_pt=16,
            color_hex="#999999",
        )


def build_ending(slide, slide_data: dict, design: dict):
    """Build a closing/thank-you slide."""
    palette = design["palette"]
    fonts = design["fonts"]
    set_slide_bg(slide, palette["primary"])

    title = slide_data.get("title", "谢谢")
    subtitle = slide_data.get("subtitle", "Q & A")

    add_textbox(
        slide, left=2.0, top=2.5, width=9.3, height=1.5,
        text=title, font_name=fonts["title"], font_size_pt=48,
        bold=True, color_hex="#ffffff", alignment=PP_ALIGN.CENTER,
    )
    add_shape_bg(slide, 5.5, 4.1, 2.3, 0.05, palette["accent"])
    add_textbox(
        slide, left=2.0, top=4.4, width=9.3, height=0.8,
        text=subtitle, font_name=fonts["body"],
        font_size_pt=24, bold=False, color_hex="#cccccc",
        alignment=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Builder dispatch
# ---------------------------------------------------------------------------
SLIDE_BUILDERS = {
    "cover": build_cover,
    "toc": build_toc,
    "content": build_content,
    "chart": build_chart,
    "table": build_table,
    "image": build_image_slide,
    "diagram": build_diagram,
    "ending": build_ending,
}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def merge_design(config: dict) -> dict:
    """Merge scenario preset with explicit design overrides."""
    meta = config.get("meta", {})
    scenario = meta.get("scenario", "general")
    preset = deepcopy(SCENARIO_PRESETS.get(scenario, SCENARIO_PRESETS["general"]))
    user_design = config.get("design", {})

    # Deep merge user overrides
    for key in ["palette", "fonts"]:
        if key in user_design:
            preset[key].update(user_design[key])
    # Slide dimensions
    if "slide_width_inches" in user_design:
        preset["slide_width"] = user_design["slide_width_inches"]
    if "slide_height_inches" in user_design:
        preset["slide_height"] = user_design["slide_height_inches"]
    # Mode & template
    preset["mode"] = user_design.get("mode", "from-scratch")
    preset["template_path"] = user_design.get("template_path")

    return preset


def generate(config_path: str, output_path: str):
    """Main entry point: read config, build PPTX, save to output."""
    with open(config_path, "r", encoding="utf-8") as f:
        config = json.load(f)

    design = merge_design(config)
    slides_data = config.get("slides", [])

    # Determine presentation source
    mode = design.get("mode", "from-scratch")
    template_path = design.get("template_path")

    if mode == "template" and template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Set slide dimensions
        prs.slide_width = emu(design.get("slide_width", 13.333))
        prs.slide_height = emu(design.get("slide_height", 7.5))

    for sd in slides_data:
        slide_type = sd.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, build_content)

        # Use blank layout for from-scratch, or pick from template
        if mode == "template" and template_path and os.path.exists(template_path):
            # Try to find a matching layout by name
            layout_name_map = {
                "cover": "Title Slide",
                "content": "Title and Content",
                "ending": "Section Header",
            }
            layout_name = layout_name_map.get(slide_type, "Title and Content")
            slide_layout = None
            for layout in prs.slide_layouts:
                if layout.name and layout_name.lower() in layout.name.lower():
                    slide_layout = layout
                    break
            if slide_layout is None:
                slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(slide_layout)

        builder(slide, sd, design)

    prs.save(output_path)
    print(f"[OK] PPT saved: {output_path}")
    print(f"   Slides: {len(slides_data)}")
    print(f"   Design: {design.get('mode', 'from-scratch')} / scenario={config.get('meta', {}).get('scenario', 'general')}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate PPTX from JSON config")
    parser.add_argument("--config", required=True, help="Path to config JSON")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()
    generate(args.config, args.output)
